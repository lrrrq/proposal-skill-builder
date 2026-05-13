"""
compiler - 案例编译（本地结构化）
"""

import json
import uuid
import shutil
from pathlib import Path
from typing import Optional, List, Dict

from .config import Config
from .db import get_connection
from .utils import now_iso
from .parser import parse_file
from .extractor import extract_fragments
from .case_manager import get_case_by_id
from .office_converter import convert_pptx_to_pdf, is_libreoffice_available


# 支持的文件格式
TEXT_EXTENSIONS = {".md", ".txt"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | IMAGE_EXTENSIONS | {".pdf", ".pptx"}


def get_source_file_by_case(case_id: str) -> Optional[Dict]:
    """根据 case_id 获取 source_file 信息"""
    conn = get_connection()
    cur = conn.cursor()
    cur.execute("""
        SELECT sf.* FROM source_files sf
        WHERE sf.case_id = ?
    """, (case_id,))
    row = cur.fetchone()
    conn.close()
    return dict(row) if row else None


def record_job(case_id: str, stage: str, status: str, error_message: str = None) -> str:
    """记录 job 到数据库"""
    job_id = f"job-{uuid.uuid4().hex[:8]}"
    now = now_iso()

    conn = get_connection()
    cur = conn.cursor()

    if status == "pending":
        cur.execute("""
            INSERT INTO jobs (job_id, case_id, stage, status, started_at)
            VALUES (?, ?, ?, ?, ?)
        """, (job_id, case_id, stage, status, now))
    else:
        cur.execute("""
            INSERT INTO jobs (job_id, case_id, stage, status, started_at, finished_at, error_message)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, (job_id, case_id, stage, status, now, now, error_message))

    conn.commit()
    conn.close()
    return job_id


def update_case_status(case_id: str, status: str):
    """更新 case 状态"""
    now = now_iso()
    conn = get_connection()
    cur = conn.cursor()
    cur.execute("""
        UPDATE cases SET status = ?, updated_at = ? WHERE case_id = ?
    """, (status, now, case_id))
    conn.commit()
    conn.close()


def ensure_dirs(case_dir: Path):
    """确保输出目录存在"""
    (case_dir / "page_images").mkdir(parents=True, exist_ok=True)
    (case_dir / "visual_assets").mkdir(parents=True, exist_ok=True)
    (case_dir / "intermediate").mkdir(parents=True, exist_ok=True)


def compile_text_case(case_id: str, file_path: Path) -> Dict:
    """
    编译纯文本文件 (.md, .txt)
    """
    pages = parse_file(file_path)
    fragments = extract_fragments(case_id, pages)
    assets = []

    return {
        "pages": pages,
        "fragments": fragments,
        "assets": assets,
        "status": "compiled",
        "partial_reason": None,
    }


def compile_image_case(case_id: str, file_path: Path, case_dir: Path) -> Dict:
    """
    编译图片文件 (.jpg, .jpeg, .png)
    """
    # 拷贝到 visual_assets
    dest_path = case_dir / "visual_assets" / file_path.name
    shutil.copy2(file_path, dest_path)

    # 读取图片信息
    try:
        from PIL import Image
        with Image.open(dest_path) as img:
            width, height = img.size
            img_format = img.format
    except Exception:
        width, height = 0, 0
        img_format = file_path.suffix.upper().lstrip(".")

    # 生成 asset 记录
    asset_id = f"asset-{uuid.uuid4().hex[:8]}"
    asset = {
        "asset_id": asset_id,
        "case_id": case_id,
        "asset_type": "image",
        "file_name": file_path.name,
        "stored_path": str(dest_path),
        "width": width,
        "height": height,
        "format": img_format,
        "needs_vision_review": True,
        "description_status": "pending",
    }

    # 图片不生成 fragments
    pages = [{
        "page_number": 1,
        "page_type": "image",
        "text_content": [],
        "raw_text": "",
        "asset_id": asset_id,
    }]
    fragments = []

    return {
        "pages": pages,
        "fragments": fragments,
        "assets": [asset],
        "status": "compiled",
        "partial_reason": None,
    }


def compile_pdf_case(case_id: str, file_path: Path, case_dir: Path) -> Dict:
    """
    编译 PDF 文件
    """
    import fitz  # PyMuPDF

    pages = []
    fragments = []
    assets = []

    try:
        doc = fitz.open(file_path)
    except Exception as e:
        raise RuntimeError(f"无法打开 PDF: {str(e)}")

    for page_num in range(len(doc)):
        page = doc[page_num]
        page_number = page_num + 1

        # 提取文本
        text = page.get_text()
        text = text.strip()
        text_length = len(text)

        # 判断是否需要视觉审查
        needs_vision_review = text_length < 30

        # 渲染页面为 PNG
        page_image_path = case_dir / "page_images" / f"page_{page_number:03d}.png"
        try:
            mat = fitz.Matrix(2.0, 2.0)  # 2x 缩放
            pix = page.get_pixmap(matrix=mat)
            pix.save(str(page_image_path))
        except Exception:
            page_image_path = None

        # 生成 page_id
        page_id = f"{case_id}_page_{page_number:04d}"

        # 生成 asset（只有真实渲染图）
        asset_id = None
        if page_image_path and page_image_path.exists():
            asset_id = f"asset-{uuid.uuid4().hex[:8]}"
            try:
                img = fitz.open(str(page_image_path))
                width = img[0].rect.width
                height = img[0].rect.height
                img.close()
            except Exception:
                width = 0
                height = 0

            asset = {
                "asset_id": asset_id,
                "case_id": case_id,
                "asset_type": "pdf_page_image",
                "page_id": page_id,
                "page_number": page_number,
                "original_path": str(file_path),
                "compiled_path": str(page_image_path),
                "stored_path": str(page_image_path),
                "filename": f"page_{page_number:03d}.png",
                "mime_type": "image/png",
                "width": width,
                "height": height,
                "format": "PNG",
                "source": "pdf_rendered_page",
                "description_status": "pending",
                "manual_description": "",
                "needs_vision_review": needs_vision_review,
                "source_text_length": text_length,
            }
            assets.append(asset)

        # 生成 page 记录
        text_lines = [l for l in text.split("\n") if l.strip()]
        has_text = text_length > 0

        # quality_flags
        quality_flags = []
        if text_length == 0:
            quality_flags.append("no_text")
        elif text_length < 30:
            quality_flags.append("low_text")
        else:
            quality_flags.append("normal")

        summary = text[:120] if text else ""

        page_record = {
            "page_id": page_id,
            "case_id": case_id,
            "page_number": page_number,
            "page_type": "text" if text_length >= 30 else "image_heavy",
            "text_content": text_lines,
            "text_length": text_length,
            "raw_text": text,
            "image_path": str(page_image_path) if page_image_path else None,
            "has_text": has_text,
            "needs_vision_review": needs_vision_review,
            "summary": summary,
            "quality_flags": quality_flags,
            "asset_id": asset_id,
        }
        pages.append(page_record)

        # 如果有足够文本，生成 fragment（带溯源字段）
        if text_length >= 30:
            fragment_id = uuid.uuid4().hex[:12]
            fragment = {
                "fragment_id": fragment_id,
                "case_id": case_id,
                "index": page_number,
                "fragment_type": "normal" if text_length >= 200 else "short",
                "raw_text": text,
                "summary": text[:120] + "..." if len(text) > 120 else text,
                "source_type": "pdf_text",
                "source_page_id": page_id,
                "source_asset_id": asset_id,
                "confidence": 0.8,
                "quality_flags": quality_flags,
            }
            fragments.append(fragment)

    doc.close()

    return {
        "pages": pages,
        "fragments": fragments,
        "assets": assets,
        "status": "compiled",
        "partial_reason": None,
    }


def compile_pptx_case(case_id: str, file_path: Path, case_dir: Path) -> Dict:
    """
    编译 PPTX 文件

    策略：
    1. 先用 python-pptx 提取文本
    2. 尝试用 LibreOffice 转 PDF
    3. 如果成功，用 PDF 渲染生成视觉资产
    4. 如果失败，只保留文本 fragments，不生成假视觉 asset
    """
    from pptx import Presentation

    pages = []
    fragments = []
    assets = []
    partial_reason = None
    final_status = "compiled"

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        raise RuntimeError(f"无法打开 PPTX: {str(e)}")

    # 获取幻灯片尺寸（英寸转像素，假设 96 DPI）
    slide_width = int(prs.slide_width.inches * 96)
    slide_height = int(prs.slide_height.inches * 96)

    # 尝试 LibreOffice 转换
    pdf_path = None
    if is_libreoffice_available():
        pdf_path = convert_pptx_to_pdf(file_path, case_dir / "intermediate")

    for slide_num, slide in enumerate(prs.slides, 1):
        # 提取文本
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        text = "\n".join(text_parts)
        needs_vision_review = len(text.strip()) < 30

        # 生成 page 记录（先创建，asset_id 待定）
        text_lines = [l for l in text.split("\n") if l.strip()]
        page_record = {
            "page_number": slide_num,
            "page_type": "text" if len(text) >= 30 else "image_heavy",
            "text_content": text_lines,
            "raw_text": text,
            "asset_id": None,  # 待填充
            "needs_vision_review": needs_vision_review,
        }
        pages.append(page_record)

        # 如果有足够文本，生成 fragment
        if len(text) >= 30:
            page_fragments = extract_fragments(case_id, [page_record])
            fragments.extend(page_fragments)

    # 如果 LibreOffice 转换成功，用 PDF 渲染生成视觉资产
    if pdf_path and pdf_path.exists():
        # 用 PDF 渲染
        import fitz
        try:
            doc = fitz.open(pdf_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                slide_num = page_num + 1

                # 渲染页面为 PNG
                page_image_path = case_dir / "page_images" / f"slide_{slide_num:03d}.png"
                try:
                    mat = fitz.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    pix.save(str(page_image_path))
                except Exception:
                    page_image_path = None

                # 为对应 page 设置 asset_id
                if slide_num <= len(pages):
                    pages[slide_num - 1]["asset_id"] = f"asset-{uuid.uuid4().hex[:8]}" if page_image_path else None

                # 只在有真实图片时创建 asset
                if page_image_path and page_image_path.exists():
                    asset_id = pages[slide_num - 1]["asset_id"]
                    asset = {
                        "asset_id": asset_id,
                        "case_id": case_id,
                        "asset_type": "pptx_slide_image",
                        "slide_number": slide_num,
                        "file_name": f"slide_{slide_num:03d}.png",
                        "stored_path": str(page_image_path),
                        "width": slide_width,
                        "height": slide_height,
                        "needs_vision_review": True,
                        "description_status": "pending",
                        "source_text_length": len(pages[slide_num - 1]["raw_text"]),
                    }
                    assets.append(asset)

            doc.close()
        except Exception:
            # PDF 渲染失败，降级
            partial_reason = "LibreOffice 转换 PDF 成功，但页面渲染失败"
            final_status = "compiled_partial"
    else:
        # 无法生成视觉资产，只保留文本
        if is_libreoffice_available():
            partial_reason = "LibreOffice 转换失败"
        else:
            partial_reason = "LibreOffice 未安装，无法生成幻灯片视觉图"
        final_status = "compiled_partial"

    return {
        "pages": pages,
        "fragments": fragments,
        "assets": assets,
        "status": final_status,
        "partial_reason": partial_reason,
    }


def compile_case(case_id: str) -> Dict:
    """
    编译案例（本地结构化）

    支持格式：
    - .md, .txt: 纯文本
    - .jpg, .jpeg, .png: 图片
    - .pdf: PDF 文档
    - .pptx: PowerPoint 文档

    输出：
    - pages.json
    - fragments.json
    - assets.json
    - page_images/ (PDF/PPTX 页面渲染图)
    - visual_assets/ (原始图片拷贝)
    - intermediate/ (PPTX 转换的 PDF)

    Returns:
        {"success": bool, "message": str, "output_dir": str, "status": str}
    """
    # 获取 case 信息
    case = get_case_by_id(case_id)
    if case is None:
        return {"success": False, "message": f"Case 不存在: {case_id}"}

    # 获取 source_file 信息
    source_file = get_source_file_by_case(case_id)
    if source_file is None:
        return {"success": False, "message": f"Case 未绑定文件: {case_id}"}

    file_path = Path(source_file["current_path"])
    ext = file_path.suffix.lower()

    # 检查文件格式
    if ext not in SUPPORTED_EXTENSIONS:
        error_msg = f"文件格式暂不支持: {ext}。支持: {', '.join(SUPPORTED_EXTENSIONS)}"
        job_id = record_job(case_id, "compile", "failed", error_msg)
        update_case_status(case_id, "failed")
        return {"success": False, "message": error_msg, "job_id": job_id}

    # 检查文件是否存在
    if not file_path.exists():
        error_msg = f"文件不存在: {file_path}"
        job_id = record_job(case_id, "compile", "failed", error_msg)
        update_case_status(case_id, "failed")
        return {"success": False, "message": error_msg}

    # 创建输出目录
    case_dir = Config.CASES_DIR / case_id
    ensure_dirs(case_dir)

    # 根据格式编译
    try:
        if ext in TEXT_EXTENSIONS:
            result = compile_text_case(case_id, file_path)
        elif ext in IMAGE_EXTENSIONS:
            result = compile_image_case(case_id, file_path, case_dir)
        elif ext == ".pdf":
            result = compile_pdf_case(case_id, file_path, case_dir)
        elif ext == ".pptx":
            result = compile_pptx_case(case_id, file_path, case_dir)
        else:
            raise ValueError(f"不支持的格式: {ext}")
    except Exception as e:
        error_msg = f"编译失败: {str(e)}"
        job_id = record_job(case_id, "compile", "failed", error_msg)
        update_case_status(case_id, "failed")
        return {"success": False, "message": error_msg, "job_id": job_id}

    # 写入 pages.json
    pages_path = case_dir / "pages.json"
    with open(pages_path, "w", encoding="utf-8") as f:
        json.dump(result["pages"], f, ensure_ascii=False, indent=2)

    # 写入 fragments.json
    fragments_path = case_dir / "fragments.json"
    with open(fragments_path, "w", encoding="utf-8") as f:
        json.dump(result["fragments"], f, ensure_ascii=False, indent=2)

    # 写入 assets.json（只有有效 asset）
    assets_path = case_dir / "assets.json"
    with open(assets_path, "w", encoding="utf-8") as f:
        json.dump(result["assets"], f, ensure_ascii=False, indent=2)

    # 生成报告
    report_path = case_dir / "extraction_report.md"
    generate_extraction_report(
        case_id,
        result["pages"],
        result["fragments"],
        result["assets"],
        report_path,
        partial_reason=result.get("partial_reason"),
        source_ext=ext
    )

    # 更新 case 状态
    update_case_status(case_id, result["status"])

    # 记录 job
    record_job(case_id, "compile", "success")

    return {
        "success": True,
        "message": f"案例编译完成（状态: {result['status']}）",
        "output_dir": str(case_dir),
        "pages_count": len(result["pages"]),
        "fragments_count": len(result["fragments"]),
        "assets_count": len(result["assets"]),
        "status": result["status"],
        "partial_reason": result.get("partial_reason"),
    }


def generate_extraction_report(case_id: str, pages: List[Dict], fragments: List[Dict],
                                assets: List[Dict], output_path: Path,
                                partial_reason: str = None,
                                source_ext: str = None) -> Path:
    """生成 extraction report"""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 统计
    needs_vision_count = sum(1 for a in assets if a.get("needs_vision_review"))
    pages_with_text = sum(1 for p in pages if p.get("raw_text", "").strip())
    pages_image_heavy = sum(1 for p in pages if p.get("page_type") == "image_heavy")

    # PDF 详细统计（只有 PDF 源才显示）
    is_pdf = source_ext == ".pdf"
    pdf_total_pages = len(pages) if is_pdf else 0
    pdf_pages_with_text = sum(1 for p in pages if p.get("raw_text", "").strip()) if is_pdf else 0
    pdf_pages_low_text = sum(1 for p in pages if 0 < len(p.get("raw_text", "").strip()) < 30) if is_pdf else 0
    pdf_pages_no_text = sum(1 for p in pages if not p.get("raw_text", "").strip()) if is_pdf else 0
    pdf_rendered_images = sum(1 for p in pages if p.get("asset_id") is not None) if is_pdf else 0

    lines = [
        "# Extraction Report",
        "",
        f"**Case ID**: {case_id}",
        f"**生成时间**: {now_iso()}",
        "",
        "## 统计",
        "",
        f"- **Pages**: {len(pages)}",
        f"- **Fragments**: {len(fragments)}",
        f"- **Assets**: {len(assets)}",
        f"- **需要视觉审查**: {needs_vision_count}",
        f"- **纯文本页面**: {pages_with_text}",
        f"- **图像为主页面**: {pages_image_heavy}",
        "",
    ]

    # PDF 详细统计（当源文件是 PDF 时才显示）
    if is_pdf:
        lines.extend([
            "## PDF 统计",
            "",
            f"- **PDF 总页数**: {pdf_total_pages}",
            f"- **有文本页数**: {pdf_pages_with_text}",
            f"- **低文本页数** (<30字): {pdf_pages_low_text}",
            f"- **无文本页数**: {pdf_pages_no_text}",
            f"- **渲染图数量**: {pdf_rendered_images}",
            "",
        ])

    if partial_reason:
        lines.extend([
            "## ⚠️ 部分编译",
            "",
            f"**原因**: {partial_reason}",
            "",
            "如需完整视觉理解，请：",
            "1. 安装 LibreOffice",
            "2. 或提供原始 PDF 版本",
            "",
        ])

    # Assets 详情
    if assets:
        lines.extend([
            "## Assets",
            "",
        ])
        for a in assets:
            asset_type = a.get("asset_type", "unknown")
            needs_review = "👁" if a.get("needs_vision_review") else "📄"
            lines.append(f"### {needs_review} {a['asset_id']} ({asset_type})")
            lines.append(f"- **类型**: {asset_type}")
            if a.get("slide_number"):
                lines.append(f"- **幻灯片**: {a['slide_number']}")
            if a.get("width") and a.get("height"):
                lines.append(f"- **尺寸**: {a['width']}x{a['height']}")
            lines.append(f"- **需要视觉审查**: {a.get('needs_vision_review', False)}")
            lines.append(f"- **描述状态**: {a.get('description_status', 'pending')}")
            lines.append("")
    else:
        lines.extend([
            "## Assets",
            "",
            "*无视觉资产（仅文本）*",
            "",
        ])

    # Fragments 详情
    if fragments:
        lines.extend([
            "## Fragments",
            "",
        ])
        for frag in fragments:
            lines.append(f"### {frag['fragment_id']} ({frag.get('fragment_type', 'unknown')})")
            lines.append(f"- **Summary**: {frag.get('summary', '')[:80]}...")
            lines.append("")

    lines.extend([
        "---",
        "",
        "*由 Proposal Skill Builder 自动生成*",
    ])

    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path